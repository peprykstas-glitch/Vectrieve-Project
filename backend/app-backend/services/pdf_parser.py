import asyncio
import io
import tempfile
from pathlib import Path
from typing import List, Optional

from models.document import Document, DocumentStatus, DocumentChunk
from core.database import get_session_factory


# ---------------------------------------------------------------------------
# Format extractors — pure synchronous CPU-bound functions
# ---------------------------------------------------------------------------

def extract_text_from_docx(file_bytes: bytes) -> str:
    import zipfile
    import xml.etree.ElementTree as ET
    with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
        if 'word/document.xml' not in z.namelist():
            raise ValueError("Not a valid DOCX file.")
        doc_xml = z.read('word/document.xml')
        root = ET.fromstring(doc_xml)
        ns = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
        paragraphs = []
        for p in root.findall('.//w:p', ns):
            p_text = [t.text for t in p.findall('.//w:t', ns) if t.text]
            if p_text:
                paragraphs.append("".join(p_text))
        return "\n".join(paragraphs)


def extract_text_from_epub(file_bytes: bytes) -> str:
    import zipfile
    import xml.etree.ElementTree as ET
    import os
    from html.parser import HTMLParser

    class EPUBTextStripper(HTMLParser):
        def __init__(self):
            super().__init__()
            self.text_parts: List[str] = []
            self.in_style_or_script = False

        def handle_starttag(self, tag, attrs):
            if tag in ('style', 'script'):
                self.in_style_or_script = True
            elif tag in ('p', 'div', 'li', 'br', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'tr'):
                if self.text_parts and not self.text_parts[-1].endswith('\n'):
                    self.text_parts.append('\n')

        def handle_endtag(self, tag):
            if tag in ('style', 'script'):
                self.in_style_or_script = False
            elif tag in ('p', 'div', 'li', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'tr'):
                if self.text_parts and not self.text_parts[-1].endswith('\n'):
                    self.text_parts.append('\n')

        def handle_data(self, data):
            if not self.in_style_or_script:
                self.text_parts.append(data)

        def get_text(self) -> str:
            return "".join(self.text_parts)

    with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
        try:
            container_xml = z.read("META-INF/container.xml")
            root = ET.fromstring(container_xml)
            ns = {'ns': 'urn:oasis:names:tc:opendocument:xmlns:container'}
            rootfile = root.find('.//ns:rootfile', ns)
            opf_path = rootfile.attrib['full-path']
        except Exception:
            opf_files = [f for f in z.namelist() if f.endswith('.opf')]
            if opf_files:
                opf_path = opf_files[0]
            else:
                raise ValueError("Not a valid EPUB: content.opf file not found.")

        opf_data = z.read(opf_path)
        opf_root = ET.fromstring(opf_data)
        base_dir = os.path.dirname(opf_path)
        ns_opf = {
            'opf': 'http://www.idpf.org/2007/opf',
            'dc': 'http://purl.org/dc/elements/1.1/'
        }
        manifest = opf_root.find('.//opf:manifest', ns_opf)
        items = {}
        if manifest is not None:
            for item in manifest.findall('.//opf:item', ns_opf):
                items[item.attrib['id']] = item.attrib['href']
        spine = opf_root.find('.//opf:spine', ns_opf)
        spine_item_ids = []
        if spine is not None:
            for itemref in spine.findall('.//opf:itemref', ns_opf):
                spine_item_ids.append(itemref.attrib['idref'])
        full_text = []
        for item_id in spine_item_ids:
            if item_id in items:
                href = items[item_id]
                path = os.path.join(base_dir, href) if base_dir else href
                path = path.replace('\\', '/')
                if path in z.namelist():
                    html_bytes = z.read(path)
                    try:
                        html_str = html_bytes.decode('utf-8')
                    except UnicodeDecodeError:
                        html_str = html_bytes.decode('latin-1', errors='ignore')
                    parser = EPUBTextStripper()
                    parser.feed(html_str)
                    page_text = parser.get_text().strip()
                    if page_text:
                        full_text.append(page_text)
        return "\n\n".join(full_text)


def extract_text_from_html(file_bytes: bytes) -> str:
    from html.parser import HTMLParser
    
    class HTMLTextStripper(HTMLParser):
        def __init__(self):
            super().__init__()
            self.text_parts: List[str] = []
            self.in_style_or_script = False

        def handle_starttag(self, tag, attrs):
            if tag in ('style', 'script', 'head'):
                self.in_style_or_script = True
            elif tag in ('p', 'div', 'li', 'br', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'tr'):
                if self.text_parts and not self.text_parts[-1].endswith('\n'):
                    self.text_parts.append('\n')

        def handle_endtag(self, tag):
            if tag in ('style', 'script', 'head'):
                self.in_style_or_script = False
            elif tag in ('p', 'div', 'li', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'tr'):
                if self.text_parts and not self.text_parts[-1].endswith('\n'):
                    self.text_parts.append('\n')

        def handle_data(self, data):
            if not self.in_style_or_script:
                self.text_parts.append(data)

        def get_text(self) -> str:
            return "".join(self.text_parts)

    html_str = None
    for encoding in ('utf-8', 'windows-1251', 'utf-16', 'latin-1'):
        try:
            html_str = file_bytes.decode(encoding)
            break
        except UnicodeDecodeError:
            continue
    if html_str is None:
        html_str = file_bytes.decode('latin-1', errors='ignore')

    parser = HTMLTextStripper()
    parser.feed(html_str)
    return parser.get_text().strip()


def _parse_file_sync(file_path: Path, filename: str) -> str:
    """
    Pure synchronous parser — reads file from disk and returns extracted text.
    Designed to be called via asyncio.to_thread() so it never blocks the event loop.
    """
    file_bytes = file_path.read_bytes()

    if filename.lower().endswith('.pdf'):
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(file_bytes))
        text = ""
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        return text

    if filename.lower().endswith('.pptx'):
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        slide_texts = []
        for i, slide in enumerate(prs.slides):
            parts = [f"--- Slide {i+1} ---"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            parts.append(paragraph.text.strip())
                if shape.has_table:
                    table = shape.table
                    table_data = []
                    for row in table.rows:
                        row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_cells:
                            table_data.append(" | ".join(row_cells))
                    if table_data:
                        parts.append("\n".join(table_data))
            try:
                notes_slide = slide.notes_slide
                if notes_slide and notes_slide.notes_text_frame:
                    notes_text = notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        parts.append(f"Speaker Notes:\n{notes_text}")
            except Exception:
                pass
            slide_texts.append("\n".join(parts))
        return "\n\n".join(slide_texts)

    if filename.lower().endswith('.docx'):
        return extract_text_from_docx(file_bytes)

    if filename.lower().endswith('.epub'):
        return extract_text_from_epub(file_bytes)

    if filename.lower().endswith(('.html', '.htm')):
        return extract_text_from_html(file_bytes)

    # Plain text with encoding auto-detection
    for encoding in ('utf-8', 'windows-1251', 'utf-16', 'latin-1'):
        try:
            return file_bytes.decode(encoding)
        except (UnicodeDecodeError, Exception):
            continue

    raise ValueError(
        "Unsupported binary file format. Only PDF, PPTX, DOCX, EPUB, HTML, Markdown, and text files are supported."
    )


def _sample_chunks(chunks: List[str], n: int = 8) -> List[str]:
    """
    Sample n chunks distributed evenly across the full document.
    Gives a representative cross-section instead of always reading
    the first 3 chunks (which are usually the title/TOC).
    """
    total = len(chunks)
    if total <= n:
        return chunks
    # Distribute indices evenly across document: 0%, ~14%, ~28%, ..., 100%
    step = (total - 1) / (n - 1)
    indices = sorted(set(round(i * step) for i in range(n)))
    return [chunks[i] for i in indices]


# ---------------------------------------------------------------------------
# Background task — called from upload.py via BackgroundTasks
# ---------------------------------------------------------------------------

async def process_pdf_background(doc_id: int, tmp_path: Path, filename: str, user_id: int, space_id: Optional[str] = None) -> None:
    """
    Processes an uploaded file that has already been streamed to disk at tmp_path.

    Design decisions:
    - Heavy CPU-bound parsing is dispatched to a thread pool via asyncio.to_thread()
      so the event loop is NEVER blocked during file processing (Bug 1 fix).
    - The function accepts a Path, not raw bytes, to avoid loading the entire
      file into RAM (Bug 4 fix).
    - The temp file is always deleted in the finally block.
    """
    session_factory = get_session_factory()
    from services.ws_manager import manager

    async def send_ws(status_val: str, error_msg: str = None, chunk_count: int = None, file_size: int = None):
        payload = {"type": "file_status", "doc_id": doc_id, "status": status_val}
        if error_msg:
            payload["error"] = error_msg
        if chunk_count is not None:
            payload["chunk_count"] = chunk_count
        if file_size is not None:
            payload["file_size"] = file_size
        await manager.send_personal_message(payload, user_id)

    async with session_factory() as session:
        try:
            from sqlmodel import select
            stmt = select(Document).where(Document.id == doc_id)
            result = await session.execute(stmt)
            doc = result.scalar_one_or_none()
            if not doc:
                return

            doc.status = DocumentStatus.PROCESSING.value
            await session.commit()
            await send_ws(DocumentStatus.PROCESSING.value)

            # Bug 1 fix: offload ALL CPU-bound parsing to a thread pool worker.
            # The event loop remains free to serve other users during this call.
            text = await asyncio.to_thread(_parse_file_sync, tmp_path, filename)

            # Smarter chunking using Langchain's RecursiveCharacterTextSplitter
            chunks = []
            if text:
                from langchain_text_splitters import RecursiveCharacterTextSplitter
                splitter = RecursiveCharacterTextSplitter(
                    chunk_size=1000,
                    chunk_overlap=200,
                    separators=["\n\n", "\n", " ", ""]
                )
                chunks = splitter.split_text(text)

            if not chunks:
                raise ValueError("No extractable text found in file. Make sure it contains digital text (not a scanned image).")

            # Save chunks to PostgreSQL/SQLite for hybrid search
            for idx, chunk_text in enumerate(chunks):
                db_chunk = DocumentChunk(
                    document_id=doc_id,
                    user_id=user_id,
                    content=chunk_text,
                    chunk_index=idx
                )
                session.add(db_chunk)
            await session.commit()

            # Bug 5 fix: sample chunks distributed across the full document
            # instead of only reading the first 3 (which are usually title/TOC).
            try:
                from services.llm_service import llm_service
                sampled = _sample_chunks(chunks, n=8)
                summary_input = "\n\n".join(sampled)
                summary_prompt = f"""
You are Vectrieve Core, a premium business document intelligence analyzer.
Provide a highly structured, polished, and extremely concise Executive Briefing for this document in English.
Outline:
1. Document Category (e.g. Resume/CV, SLA, NDA, Corporate Guideline, FAQ, Research)
2. High-level Summary (1-2 sentences)
3. Key Takeaways or Highlighted Skills (bullet points)
4. Key Risks, Warnings, or Compliance Issues (bullet points or "None")

Format headings clearly as bold text like **Document Category:** or **Key Takeaways:**. Use standard bullet points. Keep it professional.

Document Sample:
"{summary_input}"
"""
                from models.schemas import QueryRequest, ChatMessage
                fake_req = QueryRequest(
                    messages=[ChatMessage(role="user", content=summary_prompt)],
                    thinking_mode="auditor",
                    mode="cloud"
                )
                if not llm_service.groq_client:
                    fake_req.mode = "local"

                summary_text, _ = await llm_service.generate_response(fake_req, "")
                if summary_text:
                    doc.summary = summary_text.strip()
                    session.add(doc)
                    await session.commit()
            except Exception as sum_err:
                print(f"⚠️ Failed to generate AI document summary: {sum_err}")

            # Upsert to vector DB if available
            try:
                from services.vector_service import get_vector_service
                vs = get_vector_service()
                if vs and chunks:
                    await send_ws("EMBEDDING")
                    await vs.upsert_batch(chunks, filename, user_id, space_id=space_id)
            except Exception as e:
                print(f"⚠️ Vector upsert failed: {e}")
                raise RuntimeError(f"Vector upsert failed: {e}")

            doc.status = DocumentStatus.COMPLETED.value
            doc.chunk_count = len(chunks) if chunks else 0
            await session.commit()
            await send_ws(DocumentStatus.COMPLETED.value, chunk_count=doc.chunk_count, file_size=doc.file_size)

        except Exception as e:
            await session.rollback()

            stmt = select(Document).where(Document.id == doc_id)
            result = await session.execute(stmt)
            doc = result.scalar_one_or_none()
            if doc:
                doc.status = DocumentStatus.FAILED.value
                doc.error_log = str(e)
                await session.commit()
                await send_ws(DocumentStatus.FAILED.value, str(e))
        finally:
            # Bug 4 fix: always clean up the temp file after processing finishes
            try:
                if tmp_path.exists():
                    tmp_path.unlink()
            except Exception:
                pass