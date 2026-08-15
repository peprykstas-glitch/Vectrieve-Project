import asyncio
import io
import tempfile
from pathlib import Path
from typing import List, Optional, Union

from models.document import Document, DocumentStatus, DocumentChunk
from core.database import get_session_factory


# ---------------------------------------------------------------------------
# Format extractors — pure synchronous CPU-bound functions
# ---------------------------------------------------------------------------

def extract_text_from_docx(file_bytes: bytes) -> str:
    import zipfile
    import xml.etree.ElementTree as ET
    with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
        if 'word/document.xml' not in z.namelist():
            raise ValueError("Not a valid DOCX file.")
        doc_xml = z.read('word/document.xml')
        root = ET.fromstring(doc_xml)
        ns = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
        paragraphs = []
        for p in root.findall('.//w:p', ns):
            p_text = [t.text for t in p.findall('.//w:t', ns) if t.text]
            if p_text:
                paragraphs.append("".join(p_text))
        return "\n".join(paragraphs)


def extract_text_from_epub(file_bytes: bytes) -> str:
    import zipfile
    import xml.etree.ElementTree as ET
    import os
    from html.parser import HTMLParser

    class EPUBTextStripper(HTMLParser):
        def __init__(self):
            super().__init__()
            self.text_parts: List[str] = []
            self.in_style_or_script = False

        def handle_starttag(self, tag, attrs):
            if tag in ('style', 'script'):
                self.in_style_or_script = True
            elif tag in ('p', 'div', 'li', 'br', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'tr'):
                if self.text_parts and not self.text_parts[-1].endswith('\n'):
                    self.text_parts.append('\n')

        def handle_endtag(self, tag):
            if tag in ('style', 'script'):
                self.in_style_or_script = False
            elif tag in ('p', 'div', 'li', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'tr'):
                if self.text_parts and not self.text_parts[-1].endswith('\n'):
                    self.text_parts.append('\n')

        def handle_data(self, data):
            if not self.in_style_or_script:
                self.text_parts.append(data)

        def get_text(self) -> str:
            return "".join(self.text_parts)

    with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
        try:
            container_xml = z.read("META-INF/container.xml")
            root = ET.fromstring(container_xml)
            ns = {'ns': 'urn:oasis:names:tc:opendocument:xmlns:container'}
            rootfile = root.find('.//ns:rootfile', ns)
            opf_path = rootfile.attrib['full-path']
        except Exception:
            opf_files = [f for f in z.namelist() if f.endswith('.opf')]
            if opf_files:
                opf_path = opf_files[0]
            else:
                raise ValueError("Not a valid EPUB: content.opf file not found.")

        opf_data = z.read(opf_path)
        opf_root = ET.fromstring(opf_data)
        base_dir = os.path.dirname(opf_path)
        ns_opf = {
            'opf': 'http://www.idpf.org/2007/opf',
            'dc': 'http://purl.org/dc/elements/1.1/'
        }
        manifest = opf_root.find('.//opf:manifest', ns_opf)
        items = {}
        if manifest is not None:
            for item in manifest.findall('.//opf:item', ns_opf):
                items[item.attrib['id']] = item.attrib['href']
        spine = opf_root.find('.//opf:spine', ns_opf)
        spine_item_ids = []
        if spine is not None:
            for itemref in spine.findall('.//opf:itemref', ns_opf):
                spine_item_ids.append(itemref.attrib['idref'])
        full_text = []
        for item_id in spine_item_ids:
            if item_id in items:
                href = items[item_id]
                path = os.path.join(base_dir, href) if base_dir else href
                path = path.replace('\\', '/')
                if path in z.namelist():
                    html_bytes = z.read(path)
                    try:
                        html_str = html_bytes.decode('utf-8')
                    except UnicodeDecodeError:
                        html_str = html_bytes.decode('latin-1', errors='ignore')
                    parser = EPUBTextStripper()
                    parser.feed(html_str)
                    page_text = parser.get_text().strip()
                    if page_text:
                        full_text.append(page_text)
        return "\n\n".join(full_text)


def extract_text_from_html(file_bytes: bytes) -> str:
    from html.parser import HTMLParser
    
    class HTMLTextStripper(HTMLParser):
        def __init__(self):
            super().__init__()
            self.text_parts: List[str] = []
            self.in_style_or_script = False

        def handle_starttag(self, tag, attrs):
            if tag in ('style', 'script', 'head'):
                self.in_style_or_script = True
            elif tag in ('p', 'div', 'li', 'br', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'tr'):
                if self.text_parts and not self.text_parts[-1].endswith('\n'):
                    self.text_parts.append('\n')

        def handle_endtag(self, tag):
            if tag in ('style', 'script', 'head'):
                self.in_style_or_script = False
            elif tag in ('p', 'div', 'li', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'tr'):
                if self.text_parts and not self.text_parts[-1].endswith('\n'):
                    self.text_parts.append('\n')

        def handle_data(self, data):
            if not self.in_style_or_script:
                self.text_parts.append(data)

        def get_text(self) -> str:
            return "".join(self.text_parts)

    html_str = None
    for encoding in ('utf-8', 'windows-1251', 'utf-16', 'latin-1'):
        try:
            html_str = file_bytes.decode(encoding)
            break
        except UnicodeDecodeError:
            continue
    if html_str is None:
        html_str = file_bytes.decode('latin-1', errors='ignore')

    parser = HTMLTextStripper()
    parser.feed(html_str)
    return parser.get_text().strip()


def _group_rows_with_char_budget(
    rows: List[str], header_prefix: str, rows_per_group: int = 10, max_chars: int = 2500
) -> List[str]:
    """
    Groups formatted row strings into chunks, respecting BOTH a row-count cap
    and a hard character budget — a single verbose row (long text cell) must
    not be allowed to blow the embedding model's context window.
    """
    groups: List[str] = []
    current: List[str] = []
    current_len = len(header_prefix)

    for row in rows:
        row_len = len(row) + 1  # count newline
        would_exceed = current_len + row_len > max_chars
        at_row_cap = len(current) >= rows_per_group
        if current and (would_exceed or at_row_cap):
            groups.append(header_prefix + "\n".join(current))
            current = []
            current_len = len(header_prefix)
        current.append(row)
        current_len += row_len

    if current:
        groups.append(header_prefix + "\n".join(current))
    return groups


def parse_csv_to_chunks(file_path: Path, filename: str) -> List[str]:
    import csv

    # Detect encoding safely (no try/except wrapper — errors propagate to
    # process_pdf_background which converts them to status=FAILED + error_log).
    file_bytes = file_path.read_bytes()
    csv_str = None
    for encoding in ('utf-8', 'windows-1251', 'utf-16', 'latin-1'):
        try:
            csv_str = file_bytes.decode(encoding)
            break
        except UnicodeDecodeError:
            continue
    if csv_str is None:
        csv_str = file_bytes.decode('latin-1', errors='ignore')

    # Use io.StringIO so that csv.reader sees the full raw string and can
    # correctly handle quoted fields that contain embedded newlines, e.g.:
    #   "Great guest.\nLeft a 5-star review."
    # splitlines() would have destroyed that structure before csv.reader
    # had a chance to process the quotes.
    reader = csv.reader(io.StringIO(csv_str))
    rows = list(reader)

    if not rows:
        return []

    headers = [col.strip() for col in rows[0]]
    data_rows = rows[1:]

    if not data_rows:
        return [f"Table: {filename}\nColumns: {' | '.join(headers)}\nNo data rows."]

    header_prefix = f"=== Source File: {filename} ===\nFormat: Table Row Group\nColumns: {' | '.join(headers)}\n---\n"

    formatted_rows = []
    for idx, row in enumerate(data_rows):
        row_num = idx + 1
        row_str = " | ".join(val.strip() for val in row)
        formatted_rows.append(f"Row {row_num}: {row_str}")

    return _group_rows_with_char_budget(formatted_rows, header_prefix, rows_per_group=10, max_chars=2500)


def parse_excel_to_chunks(file_path: Path, filename: str) -> List[str]:
    import pandas as pd

    # No try/except wrapper — openpyxl errors propagate to process_pdf_background
    # which converts them to status=FAILED + error_log (avoids fake-success chunks).
    chunks = []
    dict_dfs = pd.read_excel(file_path, sheet_name=None, engine='openpyxl')

    for sheet_name, df in dict_dfs.items():
        df = df.dropna(how='all')
        if df.empty:
            continue
            
        headers = [str(col).strip() for col in df.columns]
        data_rows = df.values.tolist()
        
        if not data_rows:
            chunks.append(f"Table: {filename} (Sheet: {sheet_name})\nColumns: {' | '.join(headers)}\nNo data rows.")
            continue
            
        header_prefix = f"=== Source File: {filename} (Sheet: {sheet_name}) ===\nFormat: Table Row Group\nColumns: {' | '.join(headers)}\n---\n"
        
        formatted_rows = []
        for idx, row in enumerate(data_rows):
            row_num = idx + 1
            row_cells = []
            for val in row:
                if pd.isna(val):
                    row_cells.append("")
                else:
                    row_cells.append(str(val).strip())
            row_str = " | ".join(row_cells)
            formatted_rows.append(f"Row {row_num}: {row_str}")

        sheet_chunks = _group_rows_with_char_budget(formatted_rows, header_prefix, rows_per_group=10, max_chars=2500)
        chunks.extend(sheet_chunks)
            
    return chunks


def parse_json_to_chunks(file_path: Path, filename: str) -> List[str]:
    import json

    # Multi-encoding decode with fallback for Windows-1251 (Cyrillic) and Latin-1
    file_bytes = file_path.read_bytes()
    decoded_text = None
    for enc in ("utf-8", "utf-8-sig", "windows-1251", "latin-1"):
        try:
            decoded_text = file_bytes.decode(enc)
            break
        except UnicodeDecodeError:
            continue

    if decoded_text is None:
        decoded_text = file_bytes.decode("utf-8", errors="ignore")

    json_data = json.loads(decoded_text)

    chunks = []
    
    if isinstance(json_data, list) and all(isinstance(item, dict) for item in json_data if item):
        keys_set = []
        for item in json_data:
            for k in item.keys():
                if k not in keys_set:
                    keys_set.append(k)
        
        headers = keys_set
        header_prefix = f"=== Source File: {filename} ===\nFormat: JSON Record Group\nColumns: {' | '.join(headers)}\n---\n"
        
        formatted_rows = []
        for idx, item in enumerate(json_data):
            record_num = idx + 1
            row_cells = [str(item.get(k, "")).strip() for k in headers]
            row_str = " | ".join(row_cells)
            formatted_rows.append(f"Record {record_num}: {row_str}")
            
        list_chunks = _group_rows_with_char_budget(formatted_rows, header_prefix, rows_per_group=10, max_chars=2500)
        chunks.extend(list_chunks)
            
    elif isinstance(json_data, dict):
        for key, val in json_data.items():
            val_str = json.dumps(val, indent=2, ensure_ascii=False)
            if len(val_str) < 1500:
                chunks.append(
                    f"=== Source File: {filename} ===\n"
                    f"Key: {key}\n"
                    f"Value:\n{val_str}"
                )
            else:
                from langchain_text_splitters import RecursiveCharacterTextSplitter
                splitter = RecursiveCharacterTextSplitter(
                    chunk_size=1000,
                    chunk_overlap=200,
                    separators=["\n\n", "\n", " ", ""]
                )
                sub_chunks = splitter.split_text(val_str)
                for s_idx, sc in enumerate(sub_chunks):
                    chunks.append(
                        f"=== Source File: {filename} ===\n"
                        f"Key: {key} (Part {s_idx+1}/{len(sub_chunks)})\n"
                        f"Value:\n{sc}"
                    )
    else:
        val_str = json.dumps(json_data, indent=2, ensure_ascii=False)
        from langchain_text_splitters import RecursiveCharacterTextSplitter
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            separators=["\n\n", "\n", " ", ""]
        )
        sub_chunks = splitter.split_text(val_str)
        for s_idx, sc in enumerate(sub_chunks):
            chunks.append(
                f"=== Source File: {filename} ===\n"
                f"JSON Data (Part {s_idx+1}/{len(sub_chunks)}):\n{sc}"
            )

    return chunks


def _resolve_vision_provider(llm_svc, space) -> tuple:
    """
    Determine the vision inference provider (cloud/local) from Space config.

    Intentionally returns model_name=None — never forwards space.llm_model to
    vision calls. space.llm_model is a text chat model (e.g. 'llama-3.3-70b-
    versatile') configured for conversational Q&A. Text models reject image
    payloads; forwarding the name causes a silent FAILED in describe_image_bytes.

    The vision model name is left as None so describe_image/_run_cloud_vision/
    _run_local_vision use their own vision-specific defaults:
      - cloud: 'meta-llama/llama-4-scout-17b-16e-instruct' (Groq)
      - local: 'llava' (Ollama)

    Only the provider (cloud/local) hard-limit is respected — that IS about
    privacy/locality and must not be bypassed.

    Returns:
        (mode: str, model_name: None)
    """
    mode = "local" if not llm_svc.groq_client else "cloud"
    if space and space.llm_provider:
        mode = space.llm_provider
    return mode, None  # model intentionally None — see docstring


def _extract_images_from_shape(shape, images_list: list) -> None:
    """
    Recursively traverse presentation shapes to find MSO_SHAPE_TYPE.PICTURE elements,
    extracting their raw image blobs. Recursion is required to extract pictures
    nested inside group shapes (MSO_SHAPE_TYPE.GROUP).
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            if shape.image and shape.image.blob:
                images_list.append(shape.image.blob)
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                _extract_images_from_shape(child, images_list)
    except Exception:
        pass


def _parse_file_sync(file_path: Path, filename: str) -> Union[str, List[str], "VisionPayload", "ScannedPDFPayload", "PPTXPayload"]:
    """
    Pure synchronous parser — reads file from disk and returns extracted text,
    a List[str] of pre-built chunks, or a sentinel object for formats that
    require async vision-LLM processing (images, scanned PDFs, PPTX with pictures).

    Sentinels (VisionPayload, ScannedPDFPayload, PPTXPayload) are handled by
    process_pdf_background in the async context where llm_service is accessible.
    """
    from services.vision_pipeline import VisionPayload, ScannedPDFPayload, PPTXPayload, PPTXSlidePayload
    file_bytes = file_path.read_bytes()

    if filename.lower().endswith('.pdf'):
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(file_bytes))
        text = ""
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        if text.strip():
            return text
        # Scanned / image-only PDF: no extractable text found.
        # Return a sentinel so process_pdf_background can run the vision pipeline.
        # This is the fallback path for documents where pypdf yields nothing —
        # typically scans, photo PDFs, or PDFs with only embedded images.
        return ScannedPDFPayload(file_bytes=file_bytes, filename=filename)

    if filename.lower().endswith('.pptx'):
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        
        has_any_images = False
        slides_payload = []
        slide_texts = []

        for i, slide in enumerate(prs.slides):
            parts = [f"--- Slide {i+1} ---"]
            slide_images = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            parts.append(paragraph.text.strip())
                if shape.has_table:
                    table = shape.table
                    table_data = []
                    for row in table.rows:
                        row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_cells:
                            table_data.append(" | ".join(row_cells))
                    if table_data:
                        parts.append("\n".join(table_data))
                
                # Recursively extract pictures nested in the slide shape hierarchy
                _extract_images_from_shape(shape, slide_images)

            try:
                notes_slide = slide.notes_slide
                if notes_slide and notes_slide.notes_text_frame:
                    notes_text = notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        parts.append(f"Speaker Notes:\n{notes_text}")
            except Exception:
                pass

            slide_text = "\n".join(parts)
            slide_texts.append(slide_text)

            if slide_images:
                has_any_images = True
            
            slides_payload.append(
                PPTXSlidePayload(
                    slide_index=i + 1,
                    slide_text=slide_text,
                    images=slide_images
                )
            )

        if has_any_images:
            return PPTXPayload(filename=filename, slides=slides_payload)
        return "\n\n".join(slide_texts)

    if filename.lower().endswith('.docx'):
        return extract_text_from_docx(file_bytes)

    if filename.lower().endswith('.epub'):
        return extract_text_from_epub(file_bytes)

    if filename.lower().endswith(('.html', '.htm')):
        return extract_text_from_html(file_bytes)

    if filename.lower().endswith('.csv'):
        return parse_csv_to_chunks(file_path, filename)

    if filename.lower().endswith('.xlsx'):
        return parse_excel_to_chunks(file_path, filename)

    if filename.lower().endswith('.json'):
        return parse_json_to_chunks(file_path, filename)

    # Standalone image files — return sentinel for async vision processing
    if filename.lower().endswith(('.png', '.jpg', '.jpeg', '.webp')):
        return VisionPayload(file_bytes=file_bytes, filename=filename)

    # Plain text with encoding auto-detection
    for encoding in ('utf-8', 'windows-1251', 'utf-16', 'latin-1'):
        try:
            return file_bytes.decode(encoding)
        except UnicodeDecodeError:
            continue

    raise ValueError(
        "Unsupported binary file format. Only PDF, PPTX, DOCX, EPUB, HTML, CSV, XLSX, JSON, Markdown, and text files are supported."
    )


def _sample_chunks(chunks: List[str], n: int = 8, max_chars: int = 6000) -> List[str]:
    """
    Uniform sampling with a hard character-budget backstop.
    Guards against dense chunks blowing the context window even when chunk count is capped.
    """
    total = len(chunks)
    if total <= n:
        sampled = chunks
    else:
        step = (total - 1) / (n - 1)
        indices = sorted(set(round(i * step) for i in range(n)))
        sampled = [chunks[i] for i in indices]

    result: List[str] = []
    running_len = 0
    for chunk in sampled:
        if running_len + len(chunk) > max_chars and result:
            break
        result.append(chunk)
        running_len += len(chunk)
    return result


# ---------------------------------------------------------------------------
# Background task — called from upload.py via BackgroundTasks
# ---------------------------------------------------------------------------

async def process_pdf_background(doc_id: int, tmp_path: Path, filename: str, user_id: int, space_id: Optional[str] = None) -> None:
    """
    Processes an uploaded file that has already been streamed to disk at tmp_path.

    Design decisions:
    - Heavy CPU-bound parsing is dispatched to a thread pool via asyncio.to_thread()
      so the event loop is NEVER blocked during file processing (Bug 1 fix).
    - The function accepts a Path, not raw bytes, to avoid loading the entire
      file into RAM (Bug 4 fix).
    - The temp file is always deleted in the finally block.
    """
    session_factory = get_session_factory()
    from services.ws_manager import manager

    async def send_ws(status_val: str, error_msg: str = None, chunk_count: int = None, file_size: int = None):
        payload = {"type": "file_status", "doc_id": doc_id, "status": status_val}
        if error_msg:
            payload["error"] = error_msg
        if chunk_count is not None:
            payload["chunk_count"] = chunk_count
        if file_size is not None:
            payload["file_size"] = file_size
        await manager.send_personal_message(payload, user_id)

    async with session_factory() as session:
        try:
            from sqlmodel import select
            stmt = select(Document).where(Document.id == doc_id)
            result = await session.execute(stmt)
            doc = result.scalar_one_or_none()
            if not doc:
                return

            space = None
            if space_id:
                from models.sql_models import Space
                space_res = await session.execute(select(Space).where(Space.id == space_id))
                space = space_res.scalar_one_or_none()

            doc.status = DocumentStatus.PROCESSING.value
            await session.commit()
            await send_ws(DocumentStatus.PROCESSING.value)

            # Bug 1 fix: offload ALL CPU-bound parsing to a thread pool worker.
            # The event loop remains free to serve other users during this call.
            text_or_chunks = await asyncio.to_thread(_parse_file_sync, tmp_path, filename)

            # Vision pipeline — handled here in the async context because
            # _parse_file_sync (sync) cannot await llm_service.describe_image().
            from services.vision_pipeline import (
                VisionPayload, ScannedPDFPayload, PPTXPayload,
                describe_image_bytes, process_scanned_pdf,
            )
            from services.vision_pipeline import _resize_image_bytes
            from services.llm_service import llm_service as _llm_svc
            from services.llm_config_resolver import resolve_llm_config
            from models.schemas import QueryRequest, ChatMessage

            # Vision provider resolution — see _resolve_vision_provider docstring.
            # space.llm_model (text chat model) is intentionally NOT forwarded here.
            _vision_mode, _vision_model = _resolve_vision_provider(_llm_svc, space)


            chunks = []
            if isinstance(text_or_chunks, VisionPayload):
                # Standalone image: resize first (CPU-bound), then describe.
                resized_bytes = await asyncio.to_thread(
                    _resize_image_bytes, text_or_chunks.file_bytes
                )
                description = await describe_image_bytes(
                    resized_bytes, _llm_svc,
                    _vision_mode, None  # None → vision-specific model default
                )
                if description:
                    chunks = [
                        f"=== Source File: {filename} ===\n"
                        f"Visual Content Description:\n{description}"
                    ]
                else:
                    raise ValueError(
                        f"Vision pipeline returned no description for image '{filename}'. "
                        "Check that a vision-capable model is available."
                    )

            elif isinstance(text_or_chunks, ScannedPDFPayload):
                # Scanned PDF: one chunk per page (capped at MAX_OCR_PAGES).
                # process_scanned_pdf returns (chunks, truncation_warning);
                # if the file was truncated, we write the warning to doc.error_log.
                chunks, _truncation_warning = await process_scanned_pdf(
                    text_or_chunks.file_bytes, filename,
                    _llm_svc, _vision_mode, None  # None → vision-specific model default
                )
                if _truncation_warning:
                    doc.error_log = _truncation_warning
                    session.add(doc)
                    await session.commit()
                if not chunks:
                    raise ValueError(
                        f"Vision pipeline returned no content for scanned PDF '{filename}'. "
                        "Check that a vision-capable model is available."
                    )

            elif isinstance(text_or_chunks, PPTXPayload):
                # PPTX with images: process each slide, describe images, and reconstruct presentation.
                # Images are resized in worker threads. Descriptions are embedded inline.
                slide_texts = []
                failed_image_count = 0
                total_image_count = 0
                for slide in text_or_chunks.slides:
                    slide_parts = [slide.slide_text]
                    for img_idx, img_bytes in enumerate(slide.images):
                        total_image_count += 1
                        resized_bytes = await asyncio.to_thread(_resize_image_bytes, img_bytes)
                        description = await describe_image_bytes(
                            resized_bytes, _llm_svc,
                            _vision_mode, None
                        )
                        if description:
                            slide_parts.append(
                                f"\n[Slide Image {img_idx+1} Description]:\n{description}"
                            )
                        else:
                            failed_image_count += 1
                    slide_texts.append("\n".join(slide_parts))

                if failed_image_count:
                    warning = (
                        f"⚠️ {failed_image_count} of {total_image_count} images could not be "
                        f"described (vision model unavailable or inference failed) and were "
                        f"skipped. Text and remaining images were indexed normally."
                    )
                    doc.error_log = warning
                    session.add(doc)
                    await session.commit()

                combined_text = "\n\n".join(slide_texts)
                if combined_text.strip():
                    from langchain_text_splitters import RecursiveCharacterTextSplitter
                    splitter = RecursiveCharacterTextSplitter(
                        chunk_size=1000,
                        chunk_overlap=200,
                        separators=["\n\n", "\n", " ", ""]
                    )
                    chunks = splitter.split_text(combined_text)
                else:
                    raise ValueError(
                        f"Vision pipeline returned no content for presentation '{filename}'."
                    )

            elif isinstance(text_or_chunks, list):
                chunks = text_or_chunks
            elif isinstance(text_or_chunks, str) and text_or_chunks:
                from langchain_text_splitters import RecursiveCharacterTextSplitter
                splitter = RecursiveCharacterTextSplitter(
                    chunk_size=1000,
                    chunk_overlap=200,
                    separators=["\n\n", "\n", " ", ""]
                )
                chunks = splitter.split_text(text_or_chunks)

            if not chunks:
                raise ValueError("No extractable text found in file. Make sure it contains digital text (not a scanned image).")

            # Save chunks to PostgreSQL/SQLite for hybrid search
            for idx, chunk_text in enumerate(chunks):
                db_chunk = DocumentChunk(
                    document_id=doc_id,
                    user_id=user_id,
                    content=chunk_text,
                    chunk_index=idx
                )
                session.add(db_chunk)
            await session.commit()

            # Bug 5 fix: sample chunks distributed across the full document
            # instead of only reading the first 3 (which are usually title/TOC).
            try:
                # All imports already resolved in the vision block above (_llm_svc,
                # resolve_llm_config, QueryRequest, ChatMessage).
                max_chars = 6000
                if space and space.max_tokens:
                    max_chars = max(1500, min(space.max_tokens * 3, 12000))

                sampled = _sample_chunks(chunks, n=8, max_chars=max_chars)
                summary_input = "\n\n".join(sampled)
                summary_prompt = f"""
You are Vectrieve Core, a premium business document intelligence analyzer.
Provide a highly structured, polished, and extremely concise Executive Briefing for this document in English.
Outline:
1. Document Category (e.g. Resume/CV, SLA, NDA, Corporate Guideline, FAQ, Research)
2. High-level Summary (1-2 sentences)
3. Key Takeaways or Highlighted Skills (bullet points)
4. Key Risks, Warnings, or Compliance Issues (bullet points or "None")

Format headings clearly as bold text like **Document Category:** or **Key Takeaways:**. Use standard bullet points. Keep it professional.

Document Sample:
"{summary_input}"
"""
                fake_req = QueryRequest(
                    messages=[ChatMessage(role="user", content=summary_prompt)],
                    thinking_mode="auditor",
                )
                if not _llm_svc.groq_client:
                    fake_req.mode = "local"

                resolve_llm_config(fake_req, space)

                summary_text, _ = await _llm_svc.generate_response(fake_req, "")
                if summary_text:
                    doc.summary = summary_text.strip()
                    session.add(doc)
                    await session.commit()
            except Exception as sum_err:
                print(f"⚠️ Failed to generate AI document summary: {sum_err}")

            # Upsert to vector DB if available
            try:
                from services.vector_service import get_vector_service
                vs = get_vector_service()
                if vs and chunks:
                    await send_ws("EMBEDDING")
                    await vs.upsert_batch(chunks, filename, user_id, space_id=space_id)
            except Exception as e:
                print(f"⚠️ Vector upsert failed: {e}")
                raise RuntimeError(f"Vector upsert failed: {e}")

            doc.status = DocumentStatus.COMPLETED.value
            doc.chunk_count = len(chunks) if chunks else 0
            await session.commit()
            await send_ws(DocumentStatus.COMPLETED.value, chunk_count=doc.chunk_count, file_size=doc.file_size)

        except Exception as e:
            await session.rollback()

            stmt = select(Document).where(Document.id == doc_id)
            result = await session.execute(stmt)
            doc = result.scalar_one_or_none()
            if doc:
                doc.status = DocumentStatus.FAILED.value
                doc.error_log = str(e)
                await session.commit()
                await send_ws(DocumentStatus.FAILED.value, str(e))
        finally:
            # Bug 4 fix: always clean up the temp file after processing finishes
            try:
                if tmp_path.exists():
                    tmp_path.unlink()
            except Exception:
                pass