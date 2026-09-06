"""
attachment_service.py — In-memory ephemeral parser for direct chat attachments.

Processes temporary files and images attached directly in chat conversations:
- Text/Documents (PDF, DOCX, XLSX, TXT, CSV, JSON, PPTX) are parsed purely in-memory.
- Images (PNG, JPG, JPEG, WEBP, screenshots) are processed via Multimodal Vision LLMs.
- ZERO disk persistence, ZERO database rows, and ZERO Qdrant vector database pollution.
"""

from __future__ import annotations

import base64
import io
import logging
from pathlib import Path
from typing import Optional, TYPE_CHECKING, Tuple

from models.schemas import ChatAttachment

if TYPE_CHECKING:
    from services.llm_service import LLMService

logger = logging.getLogger(__name__)

# Maximum characters per ephemeral text attachment to protect prompt budget
MAX_ATTACHMENT_CHARS = 12000

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif"}
DOC_EXTENSIONS = {
    ".pdf", ".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls",
    ".txt", ".md", ".json", ".csv", ".py", ".js", ".ts", ".html",
    ".css", ".xml", ".yaml", ".yml", ".sql", ".sh", ".log"
}


def _clean_base64_data(base64_str: str) -> bytes:
    """Strip data URI prefix (e.g., 'data:image/png;base64,...') and decode bytes."""
    if "," in base64_str:
        base64_str = base64_str.split(",", 1)[1]
    return base64.b64decode(base64_str)


def _extract_pdf_in_memory(raw_bytes: bytes) -> str:
    """Extract text from PDF using pypdfium2 in-memory."""
    try:
        import pypdfium2 as pdfium
        pdf = pdfium.PdfDocument(raw_bytes)
        text_parts = []
        for i, page in enumerate(pdf):
            textpage = page.get_textpage()
            page_text = textpage.get_text_range()
            if page_text and page_text.strip():
                text_parts.append(f"[Page {i + 1}]\n{page_text.strip()}")
            if len(text_parts) >= 20:  # limit to top 20 pages for chat context
                break
        return "\n\n".join(text_parts).strip()
    except Exception as e:
        logger.warning(f"pypdfium2 extraction failed, trying pypdf fallback: {e}")
        try:
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(raw_bytes))
            parts = []
            for i, p in enumerate(reader.pages[:20]):
                txt = p.extract_text()
                if txt and txt.strip():
                    parts.append(f"[Page {i + 1}]\n{txt.strip()}")
            return "\n\n".join(parts).strip()
        except Exception as e2:
            logger.error(f"Failed to extract PDF in-memory: {e2}")
            return ""


def _extract_docx_in_memory(raw_bytes: bytes) -> str:
    """Extract text from DOCX in-memory."""
    try:
        import docx
        doc = docx.Document(io.BytesIO(raw_bytes))
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        return "\n".join(paragraphs).strip()
    except Exception as e:
        logger.error(f"Failed to extract DOCX in-memory: {e}")
        return ""


def _extract_pptx_in_memory(raw_bytes: bytes) -> str:
    """Extract text from PPTX in-memory."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(raw_bytes))
        slide_texts = []
        for i, slide in enumerate(prs.slides[:15]):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            texts.append(paragraph.text.strip())
            if texts:
                slide_texts.append(f"[Slide {i + 1}]\n" + "\n".join(texts))
        return "\n\n".join(slide_texts).strip()
    except Exception as e:
        logger.error(f"Failed to extract PPTX in-memory: {e}")
        return ""


def _extract_xlsx_in_memory(raw_bytes: bytes) -> str:
    """Extract tabular text from XLSX/XLS in-memory."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(raw_bytes), data_only=True)
        sheet_texts = []
        for sheet_name in wb.sheetnames[:3]:
            ws = wb[sheet_name]
            rows = []
            for row in list(ws.iter_rows(values_only=True))[:50]:  # top 50 rows
                clean_row = [str(c).strip() for c in row if c is not None]
                if clean_row:
                    rows.append(" | ".join(clean_row))
            if rows:
                sheet_texts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
        return "\n\n".join(sheet_texts).strip()
    except Exception as e:
        logger.error(f"Failed to extract XLSX in-memory: {e}")
        return ""


def _extract_plaintext_in_memory(raw_bytes: bytes) -> str:
    """Extract text from plain text files."""
    for enc in ("utf-8", "utf-8-sig", "cp1251", "latin-1"):
        try:
            return raw_bytes.decode(enc).strip()
        except UnicodeDecodeError:
            continue
    return ""


async def process_ephemeral_attachment(
    attachment: ChatAttachment,
    llm_service: Optional[LLMService] = None,
    groq_api_key: Optional[str] = None,
    mode: str = "cloud",
) -> Tuple[str, str]:
    """
    Process an in-chat attachment in-memory and return:
    (formatted_context_block, source_label)
    
    Guarantees:
    - Never raises unhandled exceptions (best-effort graceful fallback).
    - ZERO disk storage and ZERO vector database writes.
    """
    filename = attachment.filename or "attachment"
    ext = Path(filename).suffix.lower()

    # If text was already extracted on frontend or pre-supplied
    if attachment.extracted_text and attachment.extracted_text.strip():
        text = attachment.extracted_text.strip()[:MAX_ATTACHMENT_CHARS]
        context = (
            f"=== [DIRECT IN-CHAT ATTACHMENT: {filename}] ===\n"
            f"{text}\n"
            f"=================================================="
        )
        return context, f"📎 {filename}"

    if not attachment.base64_data:
        return "", ""

    try:
        raw_bytes = _clean_base64_data(attachment.base64_data)
    except Exception as e:
        logger.error(f"Failed to decode base64 for attachment '{filename}': {e}")
        return "", ""

    extracted_content = ""

    # 1. Image / Photo / WhatsApp Screenshot Processing
    if ext in IMAGE_EXTENSIONS or attachment.content_type.startswith("image/"):
        if llm_service:
            try:
                logger.info(f"🖼️ Running Multimodal Vision on direct chat attachment: '{filename}'...")
                vision_prompt = (
                    "You are Neurach Vision AI. Carefully analyze this image/screenshot. "
                    "Transcribe all text, numbers, dates, error messages, client requests, "
                    "or tabular items visible in the image exactly. Describe key visual details."
                )
                description = await llm_service.describe_image(
                    image_bytes=raw_bytes,
                    mode=mode,
                    vision_prompt=vision_prompt,
                    groq_api_key=groq_api_key,
                )
                if description and description.strip():
                    extracted_content = (
                        f"[Visual Transcription & Analysis of '{filename}']:\n"
                        f"{description.strip()}"
                    )
            except Exception as e:
                logger.warning(f"⚠️ Vision analysis failed for '{filename}': {e}")
                extracted_content = f"[Attached Image: '{filename}' (Visual analysis unavailable)]"
        else:
            extracted_content = f"[Attached Image: '{filename}']"

    # 2. PDF Document Processing
    elif ext == ".pdf" or attachment.content_type == "application/pdf":
        extracted_content = _extract_pdf_in_memory(raw_bytes)

    # 3. Word DOCX Processing
    elif ext == ".docx" or "wordprocessingml" in attachment.content_type:
        extracted_content = _extract_docx_in_memory(raw_bytes)

    # 4. PowerPoint PPTX Processing
    elif ext == ".pptx" or "presentationml" in attachment.content_type:
        extracted_content = _extract_pptx_in_memory(raw_bytes)

    # 5. Excel XLSX Processing
    elif ext in (".xlsx", ".xls") or "spreadsheetml" in attachment.content_type:
        extracted_content = _extract_xlsx_in_memory(raw_bytes)

    # 6. Plaintext / Code / JSON / CSV
    else:
        extracted_content = _extract_plaintext_in_memory(raw_bytes)

    if not extracted_content or not extracted_content.strip():
        logger.warning(f"No text extracted from direct attachment: '{filename}'")
        extracted_content = f"[Attached file: '{filename}' - no readable text content]"

    # Apply token budget safety cap
    if len(extracted_content) > MAX_ATTACHMENT_CHARS:
        extracted_content = extracted_content[:MAX_ATTACHMENT_CHARS] + "\n... [Content truncated to fit context budget]"

    formatted_block = (
        f"=== [DIRECT IN-CHAT ATTACHMENT: {filename}] ===\n"
        f"{extracted_content}\n"
        f"=================================================="
    )
    return formatted_block, f"📎 {filename}"
