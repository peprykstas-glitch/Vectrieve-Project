import io
import zipfile
from pathlib import Path
import pytest
from services.pdf_parser import extract_text_from_docx, extract_text_from_epub, _sample_chunks


def test_extract_text_from_docx():
    docx_bytes = io.BytesIO()
    with zipfile.ZipFile(docx_bytes, 'w') as z:
        z.writestr('word/document.xml', """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
    <w:body>
        <w:p>
            <w:r>
                <w:t>Hello world from docx paragraph 1.</w:t>
            </w:r>
        </w:p>
        <w:p>
            <w:r>
                <w:t>Second paragraph content.</w:t>
            </w:r>
        </w:p>
    </w:body>
</w:document>""")

    text = extract_text_from_docx(docx_bytes.getvalue())
    assert "Hello world from docx paragraph 1." in text
    assert "Second paragraph content." in text


def test_extract_text_from_epub():
    epub_bytes = io.BytesIO()
    with zipfile.ZipFile(epub_bytes, 'w') as z:
        z.writestr('META-INF/container.xml', """<?xml version="1.0"?>
<container version="1.0" xmlns="urn:oasis:names:tc:opendocument:xmlns:container">
  <rootfiles>
    <rootfile full-path="OEBPS/content.opf" media-type="application/oebps-package+xml"/>
  </rootfiles>
</container>""")
        z.writestr('OEBPS/content.opf', """<?xml version="1.0" encoding="UTF-8"?>
<package xmlns="http://www.idpf.org/2007/opf" unique-identifier="BookID" version="2.0">
  <metadata xmlns:dc="http://purl.org/dc/elements/1.1/">
    <dc:title>Test EPUB</dc:title>
  </metadata>
  <manifest>
    <item id="chapter1" href="chapter1.xhtml" media-type="application/xhtml+xml"/>
  </manifest>
  <spine>
    <itemref idref="chapter1"/>
  </spine>
</package>""")
        z.writestr('OEBPS/chapter1.xhtml', """<?xml version="1.0" encoding="UTF-8"?>
<html xmlns="http://www.w3.org/1999/xhtml">
  <head><title>Chapter 1</title></head>
  <body>
    <h1>Chapter 1 Header</h1>
    <p>This is the first chapter text in the EPUB.</p>
  </body>
</html>""")

    text = extract_text_from_epub(epub_bytes.getvalue())
    assert "Chapter 1 Header" in text
    assert "This is the first chapter text in the EPUB." in text


# --- Bug 5 fix: _sample_chunks tests ---

def test_sample_chunks_small():
    """When chunks <= n, return all chunks unchanged."""
    chunks = ["a", "b", "c"]
    result = _sample_chunks(chunks, n=8)
    assert result == chunks


def test_sample_chunks_distributed():
    """For large doc: samples are spread across the full document, not just the start."""
    chunks = [f"chunk_{i}" for i in range(100)]
    result = _sample_chunks(chunks, n=8)
    assert len(result) == 8
    # First and last chunks must be included (boundary coverage)
    assert result[0] == "chunk_0"
    assert result[-1] == "chunk_99"
    # Samples must NOT just be the first 8 — that was the original bug
    assert result != chunks[:8]


def test_sample_chunks_exact_n():
    """When chunks == n, all chunks returned."""
    chunks = [f"c_{i}" for i in range(8)]
    result = _sample_chunks(chunks, n=8)
    assert result == chunks


# --- Bug 4 fix: _parse_file_sync reads from Path, not bytes ---

def test_parse_file_sync_text(tmp_path):
    """_parse_file_sync correctly reads plain text files from disk."""
    from services.pdf_parser import _parse_file_sync
    p = tmp_path / "test.txt"
    p.write_text("Hello from a text file.", encoding="utf-8")
    result = _parse_file_sync(p, "test.txt")
    assert "Hello from a text file." in result


def test_parse_file_sync_docx(tmp_path):
    """_parse_file_sync correctly parses a DOCX via Path."""
    from services.pdf_parser import _parse_file_sync
    docx_bytes = io.BytesIO()
    with zipfile.ZipFile(docx_bytes, 'w') as z:
        z.writestr('word/document.xml', """<?xml version="1.0" encoding="UTF-8"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:body>
    <w:p><w:r><w:t>Docx via path test.</w:t></w:r></w:p>
  </w:body>
</w:document>""")
    p = tmp_path / "test.docx"
    p.write_bytes(docx_bytes.getvalue())
    result = _parse_file_sync(p, "test.docx")
    assert "Docx via path test." in result


def test_parse_file_sync_markdown(tmp_path):
    """_parse_file_sync correctly reads markdown files as plain text."""
    from services.pdf_parser import _parse_file_sync
    p = tmp_path / "test.md"
    p.write_text("# Markdown Title\n- Item 1\n- Item 2", encoding="utf-8")
    result = _parse_file_sync(p, "test.md")
    assert "# Markdown Title" in result
    assert "- Item 1" in result


def test_parse_file_sync_html(tmp_path):
    """_parse_file_sync correctly parses HTML files, stripping tags, head, style, and scripts."""
    from services.pdf_parser import _parse_file_sync
    # Include void elements <meta> and <link> to verify they do not cause infinite text suppression
    html_content = """
    <html>
      <head>
        <title>Ignored Title</title>
        <style>body { color: red; }</style>
        <script>console.log('ignored');</script>
      </head>
      <body>
        <meta charset="utf-8">
        <h1>Heading Content</h1>
        <link rel="stylesheet" href="style.css">
        <p>Paragraph text here.</p>
      </body>
    </html>
    """
    p = tmp_path / "test.html"
    p.write_text(html_content, encoding="utf-8")
    result = _parse_file_sync(p, "test.html")
    # Verify that clean prose text is extracted
    assert "Heading Content" in result
    assert "Paragraph text here." in result
    # Verify that headers, styles, and scripts are stripped
    assert "Ignored Title" not in result
    assert "body { color: red; }" not in result
    assert "console.log" not in result


def test_parse_file_sync_html_windows1251(tmp_path):
    """_parse_file_sync correctly decodes and parses HTML in windows-1251 encoding."""
    from services.pdf_parser import _parse_file_sync
    # Cyrillic text in windows-1251
    html_content = "<html><body><h1>Привіт Світ</h1></body></html>"
    p = tmp_path / "test_1251.html"
    p.write_bytes(html_content.encode("windows-1251"))
    result = _parse_file_sync(p, "test_1251.html")
    assert "Привіт Світ" in result


def test_parse_file_sync_csv(tmp_path):
    """_parse_file_sync correctly parses CSV into row-grouped chunks with headers."""
    from services.pdf_parser import _parse_file_sync
    csv_content = "Name,Role,Salary\nJohn,Sales,100\nJane,Support,150\nBob,Dev,200"
    p = tmp_path / "test.csv"
    p.write_text(csv_content, encoding="utf-8")
    
    chunks = _parse_file_sync(p, "test.csv")
    assert isinstance(chunks, list)
    assert len(chunks) == 1
    assert "=== Source File: test.csv ===" in chunks[0]
    assert "Columns: Name | Role | Salary" in chunks[0]
    assert "Row 1: John | Sales | 100" in chunks[0]
    assert "Row 2: Jane | Support | 150" in chunks[0]
    assert "Row 3: Bob | Dev | 200" in chunks[0]


def test_parse_file_sync_excel(tmp_path):
    """_parse_file_sync correctly parses Excel worksheets into row-grouped chunks with headers."""
    from services.pdf_parser import _parse_file_sync
    import pandas as pd
    
    df = pd.DataFrame({
        "Name": ["John", "Jane"],
        "Role": ["Sales", "Support"]
    })
    
    p = tmp_path / "test.xlsx"
    df.to_excel(p, index=False, engine="openpyxl")
    
    chunks = _parse_file_sync(p, "test.xlsx")
    assert isinstance(chunks, list)
    assert len(chunks) == 1
    assert "=== Source File: test.xlsx (Sheet: Sheet1) ===" in chunks[0]
    assert "Columns: Name | Role" in chunks[0]
    assert "Row 1: John | Sales" in chunks[0]
    assert "Row 2: Jane | Support" in chunks[0]


def test_parse_file_sync_json_list(tmp_path):
    """_parse_file_sync correctly parses flat list-of-dicts JSON into row-grouped chunks."""
    from services.pdf_parser import _parse_file_sync
    import json
    
    data = [
        {"Name": "John", "Role": "Sales"},
        {"Name": "Jane", "Role": "Support"}
    ]
    p = tmp_path / "test.json"
    p.write_text(json.dumps(data), encoding="utf-8")
    
    chunks = _parse_file_sync(p, "test.json")
    assert isinstance(chunks, list)
    assert len(chunks) == 1
    assert "Format: JSON Record Group" in chunks[0]
    assert "Columns: Name | Role" in chunks[0]
    assert "Record 1: John | Sales" in chunks[0]


def test_parse_file_sync_json_nested(tmp_path):
    """_parse_file_sync correctly chunks nested config JSON by top-level keys."""
    from services.pdf_parser import _parse_file_sync
    import json
    
    data = {
        "settings": {"theme": "dark", "version": 1},
        "users": ["user1", "user2"]
    }
    p = tmp_path / "test_nested.json"
    p.write_text(json.dumps(data), encoding="utf-8")
    
    chunks = _parse_file_sync(p, "test_nested.json")
    assert isinstance(chunks, list)
    assert len(chunks) == 2
    
    keys = [c for c in chunks if "Key: settings" in c or "Key: users" in c]
    assert len(keys) == 2


def test_group_rows_with_char_budget():
    """_group_rows_with_char_budget splits records dynamically using a character limit."""
    from services.pdf_parser import _group_rows_with_char_budget
    rows = ["Row 1: " + "A" * 1000, "Row 2: " + "B" * 1000, "Row 3: " + "C" * 1000]
    header = "Header\n---\n"
    
    result = _group_rows_with_char_budget(rows, header, rows_per_group=10, max_chars=2100)
    assert len(result) == 2
    assert result[0] == "Header\n---\nRow 1: " + "A" * 1000 + "\nRow 2: " + "B" * 1000


def test_parse_csv_multiline_field(tmp_path):
    """Regression test: quoted CSV fields containing embedded newlines must be
    kept as a single cell value, not split across multiple rows.
    The old csv.reader(html_str.splitlines()) would have broken this case."""
    from services.pdf_parser import _parse_file_sync
    # Notes column contains a newline inside a quoted field
    csv_content = (
        'Name,Role,Notes\r\n'
        'John Smith,Sales,"Great guest.\nLeft a 5-star review."\r\n'
        'Jane Doe,Marketing,OK\r\n'
    )
    p = tmp_path / "multiline.csv"
    p.write_text(csv_content, encoding="utf-8")

    chunks = _parse_file_sync(p, "multiline.csv")
    assert isinstance(chunks, list)
    # Must be exactly 2 data rows (John and Jane), not 3+ due to the embedded newline
    combined = "\n".join(chunks)
    assert "Row 1: John Smith | Sales | Great guest." in combined
    assert "Row 2: Jane Doe | Marketing | OK" in combined
    # The embedded newline continuation must NOT appear as a separate row
    assert "Row 3:" not in combined


# --- Phase 3c-core: Vision pipeline tests ---

def test_vision_mode_ignores_space_llm_model():
    """Vision calls must never inherit space.llm_model — that is a text chat model.

    Regression test for the bug where resolve_llm_config was called for _vision_req,
    propagating space.llm_model (e.g. 'llama-3.3-70b-versatile') as model_name to
    describe_image_bytes. Text models reject vision payloads, causing a silent FAILED.

    _resolve_vision_provider must return model_name=None regardless of space config.
    Only the provider (cloud/local) hard-limit is respected.
    """
    from unittest.mock import MagicMock
    from services.pdf_parser import _resolve_vision_provider

    # Space with both a custom provider AND a custom text model
    mock_space = MagicMock()
    mock_space.llm_provider = "cloud"
    mock_space.llm_model = "llama-3.3-70b-versatile"   # text model, NOT vision

    # LLM service with Groq available
    mock_svc = MagicMock()
    mock_svc.groq_client = True  # groq available → default would be "cloud"

    mode, model = _resolve_vision_provider(mock_svc, mock_space)

    assert mode == "cloud", f"Expected provider 'cloud' from space.llm_provider, got {mode!r}"
    assert model is None, (
        f"model_name must be None for vision calls — space.llm_model ({mock_space.llm_model!r}) "
        "must never be forwarded to vision. Got: {model!r}"
    )

def test_vision_provider_respects_local_hard_limit():
    """If space.llm_provider='local', vision calls must use local mode even if groq is available."""
    from unittest.mock import MagicMock
    from services.pdf_parser import _resolve_vision_provider

    mock_space = MagicMock()
    mock_space.llm_provider = "local"
    mock_space.llm_model = "some-text-model"

    mock_svc = MagicMock()
    mock_svc.groq_client = True  # groq is available, but space says local

    mode, model = _resolve_vision_provider(mock_svc, mock_space)

    assert mode == "local", f"Expected 'local' hard-limit from space, got {mode!r}"
    assert model is None


def test_render_page_to_jpeg_real_pdfium(tmp_path):
    """_render_page_to_jpeg calls real pypdfium2 C-level APIs — no mocking of pdfium.

    Creates a minimal PDF using pypdf.PdfWriter (already in deps), writes it to
    a temp file, opens with real pdfium.PdfDocument, renders page 0, and verifies
    the result is valid JPEG bytes with dimensions within MAX_IMAGE_SIDE.

    Purpose: catch API mismatches (doc[i] vs doc.get_page(i), bitmap method names,
    to_pil() signature) that are invisible to mock-only tests but crash immediately
    on the first real scanned PDF in production.
    """
    import io as _io
    import pypdfium2 as pdfium
    from PIL import Image as PILImage
    from pypdf import PdfWriter
    from services.vision_pipeline import _render_page_to_jpeg, MAX_IMAGE_SIDE

    # Build a minimal single-page PDF with pypdf
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    pdf_buf = _io.BytesIO()
    writer.write(pdf_buf)
    pdf_bytes = pdf_buf.getvalue()

    # Open with real pdfium — this exercises the actual C-level PDF parser
    doc = pdfium.PdfDocument(pdf_bytes)
    try:
        assert len(doc) == 1, "Expected exactly 1 page in test PDF"

        jpeg_bytes = _render_page_to_jpeg(doc, 0)
    finally:
        doc.close()

    # Result must be non-empty valid JPEG
    assert len(jpeg_bytes) > 0, "render returned empty bytes"
    assert jpeg_bytes[:2] == b'\xff\xd8', "Result is not a valid JPEG (wrong magic bytes)"

    # Dimensions must be within resize cap
    img = PILImage.open(_io.BytesIO(jpeg_bytes))
    w, h = img.size
    assert max(w, h) <= MAX_IMAGE_SIDE, (
        f"Rendered page exceeds MAX_IMAGE_SIDE={MAX_IMAGE_SIDE}: got {w}×{h}"
    )


def test_image_sentinel_roundtrip(tmp_path):
    """_parse_file_sync returns a VisionPayload sentinel for image files,
    NOT a string or list — confirming vision processing is deferred to async context."""
    from services.pdf_parser import _parse_file_sync
    from services.vision_pipeline import VisionPayload

    # Minimal 1x1 white JPEG (valid image bytes)
    import struct, zlib
    # Use a tiny valid JPEG (smallest possible)
    tiny_jpeg = (
        b'\xff\xd8\xff\xe0\x00\x10JFIF\x00\x01\x01\x00\x00\x01\x00\x01\x00\x00'
        b'\xff\xdb\x00C\x00\x08\x06\x06\x07\x06\x05\x08\x07\x07\x07\t\t'
        b'\x08\n\x0c\x14\r\x0c\x0b\x0b\x0c\x19\x12\x13\x0f\x14\x1d\x1a'
        b'\x1f\x1e\x1d\x1a\x1c\x1c $.\' ",#\x1c\x1c(7),01444\x1f\'9=82<.342\x1e>'
        b'\xff\xc0\x00\x0b\x08\x00\x01\x00\x01\x01\x01\x11\x00'
        b'\xff\xc4\x00\x1f\x00\x00\x01\x05\x01\x01\x01\x01\x01\x01\x00\x00'
        b'\x00\x00\x00\x00\x00\x00\x01\x02\x03\x04\x05\x06\x07\x08\t\n\x0b'
        b'\xff\xc4\x00\xb5\x10\x00\x02\x01\x03\x03\x02\x04\x03\x05\x05\x04'
        b'\x04\x00\x00\x01}\x01\x02\x03\x00\x04\x11\x05\x12!1A\x06\x13Qa'
        b'\x07"q\x142\x81\x91\xa1\x08#B\xb1\xc1\x15R\xd1\xf0$3br\x82\t\n'
        b'\xff\xda\x00\x08\x01\x01\x00\x00?\x00\xf5\x00\xff\xd9'
    )
    p = tmp_path / "test_image.jpg"
    p.write_bytes(tiny_jpeg)

    result = _parse_file_sync(p, "test_image.jpg")
    assert isinstance(result, VisionPayload), (
        f"Expected VisionPayload sentinel, got {type(result).__name__}"
    )
    assert result.filename == "test_image.jpg"
    assert result.file_bytes == tiny_jpeg


def test_png_sentinel_roundtrip(tmp_path):
    """_parse_file_sync returns VisionPayload for .png files too."""
    from services.pdf_parser import _parse_file_sync
    from services.vision_pipeline import VisionPayload

    # Minimal 1x1 red PNG
    import zlib, struct
    def create_png():
        def chunk(name, data):
            c = name + data
            return struct.pack('>I', len(data)) + c + struct.pack('>I', zlib.crc32(c) & 0xffffffff)
        sig = b'\x89PNG\r\n\x1a\n'
        ihdr = chunk(b'IHDR', struct.pack('>IIBBBBB', 1, 1, 8, 2, 0, 0, 0))
        raw = b'\x00\xff\x00\x00'  # filter byte + RGB
        idat = chunk(b'IDAT', zlib.compress(raw))
        iend = chunk(b'IEND', b'')
        return sig + ihdr + idat + iend

    p = tmp_path / "test_image.png"
    p.write_bytes(create_png())

    result = _parse_file_sync(p, "test_image.png")
    assert isinstance(result, VisionPayload)


def test_describe_image_bytes_best_effort_on_failure():
    """describe_image_bytes returns '' instead of raising when vision LLM fails.
    This ensures a single bad image never crashes the whole document ingestion."""
    import asyncio
    from unittest.mock import AsyncMock, MagicMock
    from services.vision_pipeline import describe_image_bytes

    mock_llm = MagicMock()
    mock_llm.describe_image = AsyncMock(side_effect=RuntimeError("No vision model available"))

    result = asyncio.run(
        describe_image_bytes(b"fake_image_bytes", mock_llm, mode="local", model_name=None)
    )
    assert result == "", f"Expected empty string on failure, got: {repr(result)}"



def test_process_scanned_pdf_page_limit(tmp_path):
    """process_scanned_pdf truncates to MAX_OCR_PAGES when PDF has more pages.

    Mocks _render_page_to_jpeg (the sync helper now called via asyncio.to_thread)
    rather than pdfium.PdfDocument directly, reflecting the refactored architecture
    where CPU-bound rasterization is isolated from the async loop.
    """
    import asyncio
    from unittest.mock import AsyncMock, MagicMock, patch
    from services.vision_pipeline import process_scanned_pdf, MAX_OCR_PAGES

    mock_llm = MagicMock()
    mock_llm.describe_image = AsyncMock(return_value="Page description text.")

    mock_doc = MagicMock()
    mock_doc.__len__ = MagicMock(return_value=35)  # 35 > MAX_OCR_PAGES=30
    mock_doc.close = MagicMock()

    with patch("services.vision_pipeline.pdfium") as mock_pdfium, \
         patch("services.vision_pipeline._render_page_to_jpeg", return_value=b"fake_jpeg") as mock_render:
        mock_pdfium.PdfDocument.return_value = mock_doc
        chunks, truncation_warning = asyncio.run(
            process_scanned_pdf(b"fake_pdf", "scan.pdf", mock_llm, "cloud", None)
        )

    # Only MAX_OCR_PAGES pages rendered, not 35
    assert mock_render.call_count == MAX_OCR_PAGES, (
        f"Expected {MAX_OCR_PAGES} render calls, got {mock_render.call_count}"
    )
    assert len(chunks) <= MAX_OCR_PAGES
    for chunk in chunks:
        assert "Source File: scan.pdf" in chunk
    # Truncation warning must be non-None and mention page counts
    assert truncation_warning is not None
    assert "30 of 35" in truncation_warning


def test_scanned_pdf_sentinel_fallback(tmp_path):
    """_parse_file_sync returns ScannedPDFPayload for a PDF with no extractable text.

    This is the critical path that makes the vision pipeline reachable for scanned
    documents. Previously _parse_file_sync returned an empty string for such PDFs,
    which fell through to the 'No extractable text' ValueError — vision was unreachable.
    """
    from unittest.mock import MagicMock, patch
    from services.pdf_parser import _parse_file_sync
    from services.vision_pipeline import ScannedPDFPayload

    # Simulate a scanned PDF where every page.extract_text() returns ""
    mock_page = MagicMock()
    mock_page.extract_text.return_value = ""

    mock_reader = MagicMock()
    mock_reader.pages = [mock_page, mock_page]  # two empty pages

    p = tmp_path / "scan.pdf"
    p.write_bytes(b"%PDF-1.4 fake")  # content irrelevant — reader is mocked

    # _parse_file_sync executes `from pypdf import PdfReader` inside the function.
    # Patching the PdfReader class at its source (pypdf.PdfReader) is the correct
    # target — this is what the local `from ... import` binds at call time.
    with patch("pypdf.PdfReader", return_value=mock_reader):
        result = _parse_file_sync(p, "scan.pdf")

    assert isinstance(result, ScannedPDFPayload), (
        f"Expected ScannedPDFPayload for empty PDF, got {type(result).__name__}"
    )
    assert result.filename == "scan.pdf"


def test_resize_image_bytes_caps_dimensions():
    """_resize_image_bytes caps the longest side to MAX_IMAGE_SIDE pixels.

    Verifies that a large image (3000×2000) is resized down and the result
    is a valid JPEG with longest side <= MAX_IMAGE_SIDE, not the original size.
    """
    import struct, zlib
    from PIL import Image as PILImage
    import io as _io
    from services.vision_pipeline import _resize_image_bytes, MAX_IMAGE_SIDE

    # Create a large synthetic RGB image (3000×2000) — well above MAX_IMAGE_SIDE
    large_img = PILImage.new("RGB", (3000, 2000), color=(128, 64, 32))
    buf = _io.BytesIO()
    large_img.save(buf, format="JPEG", quality=90)
    large_bytes = buf.getvalue()

    resized = _resize_image_bytes(large_bytes)

    # Load result and check dimensions
    result_img = PILImage.open(_io.BytesIO(resized))
    w, h = result_img.size
    assert max(w, h) <= MAX_IMAGE_SIDE, (
        f"Longest side {max(w, h)} exceeds MAX_IMAGE_SIDE={MAX_IMAGE_SIDE}"
    )
    # Aspect ratio preserved: 3000/2000 = 1.5 → 1920/1280 = 1.5
    assert abs(w / h - 3000 / 2000) < 0.01, "Aspect ratio was not preserved"


def test_resize_image_bytes_does_not_upscale():
    """_resize_image_bytes never enlarges a small image — thumbnail() contract."""
    from PIL import Image as PILImage
    import io as _io
    from services.vision_pipeline import _resize_image_bytes, MAX_IMAGE_SIDE

    # Small image well below MAX_IMAGE_SIDE
    small_img = PILImage.new("RGB", (400, 300), color=(0, 128, 255))
    buf = _io.BytesIO()
    small_img.save(buf, format="JPEG", quality=90)
    small_bytes = buf.getvalue()

    resized = _resize_image_bytes(small_bytes)
    result_img = PILImage.open(_io.BytesIO(resized))
    w, h = result_img.size
    # Must not be larger than original
    assert w <= 400 and h <= 300, (
        f"Image was upscaled: {w}×{h} > original 400×300"
    )


def test_parse_pptx_without_images_returns_str(tmp_path):
    """If a presentation contains only text and no images, it must bypass the vision sentinel
    and return raw text directly to minimize ingestion overhead."""
    from pptx import Presentation
    from services.pdf_parser import _parse_file_sync

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Hello Test Presentation"
    subtitle = slide.placeholders[1]
    subtitle.text = "Slide subtitle text without images"

    p = tmp_path / "text_only.pptx"
    prs.save(p)

    result = _parse_file_sync(p, "text_only.pptx")
    assert isinstance(result, str), f"Expected str for text-only presentation, got {type(result).__name__}"
    assert "Hello Test Presentation" in result
    assert "Slide subtitle text without images" in result


def test_parse_pptx_with_images_returns_payload(tmp_path):
    """If a presentation contains one or more Picture shapes (even nested inside Groups),
    it must return a PPTXPayload sentinel containing slide text and raw image bytes."""
    from pptx import Presentation
    from services.pdf_parser import _parse_file_sync
    from services.vision_pipeline import PPTXPayload
    from PIL import Image as PILImage
    import io

    # Generate a tiny PNG image
    img = PILImage.new("RGB", (10, 10), color=(255, 0, 0))
    img_buf = io.BytesIO()
    img.save(img_buf, format="PNG")
    img_bytes = img_buf.getvalue()

    prs = Presentation()
    # Slide 1: only text
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Slide 1 Title"

    # Slide 2: text + image
    slide2 = prs.slides.add_slide(prs.slide_layouts[6]) # blank
    txBox = slide2.shapes.add_textbox(0, 0, 100000, 100000)
    tf = txBox.text_frame
    tf.text = "Slide 2 text with picture"
    
    # Add image to Slide 2
    slide2.shapes.add_picture(io.BytesIO(img_bytes), 0, 200000, width=50000, height=50000)

    p = tmp_path / "with_image.pptx"
    prs.save(p)

    result = _parse_file_sync(p, "with_image.pptx")
    assert isinstance(result, PPTXPayload), f"Expected PPTXPayload, got {type(result).__name__}"
    assert len(result.slides) == 2
    
    # Slide 1 payload check
    assert result.slides[0].slide_index == 1
    assert "Slide 1 Title" in result.slides[0].slide_text
    assert len(result.slides[0].images) == 0

    # Slide 2 payload check
    assert result.slides[1].slide_index == 2
    assert "Slide 2 text with picture" in result.slides[1].slide_text
    assert len(result.slides[1].images) == 1
    assert len(result.slides[1].images[0]) > 0


def test_pptx_vision_pipeline_integration(tmp_path):
    """Verifies that PPTXPayload sentinel is correctly handled inside process_pdf_background.
    It should invoke vision descriptions for images and insert them into the reconstructed slide text
    before chunking with RecursiveCharacterTextSplitter."""
    import asyncio
    from unittest.mock import AsyncMock, MagicMock, patch
    from services.vision_pipeline import PPTXPayload, PPTXSlidePayload
    from services.pdf_parser import process_pdf_background
    from models.document import Document

    from PIL import Image as PILImage
    import io
    img = PILImage.new("RGB", (10, 10), color=(0, 255, 0))
    img_buf = io.BytesIO()
    img.save(img_buf, format="PNG")
    tiny_png_bytes = img_buf.getvalue()

    # Create dummy database document object
    doc = Document(
        id=999,
        filename="pres.pptx",
        space_id="test_space",
        status="PENDING",
        error_log=None
    )

    # Mock PPTXPayload returned by _parse_file_sync
    # Slide 1 contains 1 image blob
    slide_payload = PPTXSlidePayload(
        slide_index=1,
        slide_text="=== Slide 1 ===\nSlide text description.",
        images=[tiny_png_bytes]
    )
    payload = PPTXPayload(filename="pres.pptx", slides=[slide_payload])

    # Mocks for database, WS manager and LLM
    mock_db_session = MagicMock()
    mock_db_session.__aenter__ = AsyncMock(return_value=mock_db_session)
    mock_db_session.__aexit__ = AsyncMock()
    mock_db_session.execute = AsyncMock()
    mock_db_session.commit = AsyncMock()
    mock_db_session.add = MagicMock()

    # Stub executing select statement to return our mock doc
    mock_result = MagicMock()
    mock_result.scalar_one_or_none.return_value = doc
    mock_db_session.execute.return_value = mock_result

    # Mock WS sending
    mock_ws_manager = MagicMock()
    mock_ws_manager.send_personal_message = AsyncMock()

    # Mock LLM describes image and generates summary
    mock_llm = MagicMock()
    mock_llm.describe_image = AsyncMock(return_value="A red circle chart.")
    mock_llm.generate_response = AsyncMock(return_value=("Brief Summary", "model"))

    mock_session_factory = MagicMock(return_value=mock_db_session)

    # Mock vector service
    mock_vector_service = MagicMock()
    mock_vector_service.upsert_batch = AsyncMock()

    # Create dummy temp file so tmp_path.unlink() doesn't fail or raise warnings
    temp_file = tmp_path / "pres.pptx"
    temp_file.write_bytes(b"dummy presentation data")

    with patch("services.pdf_parser.get_session_factory", return_value=mock_session_factory), \
         patch("services.ws_manager.manager", mock_ws_manager), \
         patch("services.pdf_parser._parse_file_sync", return_value=payload), \
         patch("services.llm_service.llm_service", mock_llm), \
         patch("services.vector_service.get_vector_service", return_value=mock_vector_service), \
         patch("services.pdf_parser._resolve_vision_provider", return_value=("cloud", None)):

        # Run process_pdf_background async orchestration task
        asyncio.run(process_pdf_background(
            doc_id=999,
            tmp_path=temp_file,
            filename="pres.pptx",
            space_id="test-space",
            user_id="user1"
        ))

    # Assertions
    assert doc.status == "COMPLETED"
    assert doc.error_log is None
    
    # Verify that mock_db_session.add was called to insert chunks
    # DocumentChunk has content attribute containing chunk text
    added_chunks = []
    for call_args in mock_db_session.add.call_args_list:
        obj = call_args[0][0]
        from models.document import DocumentChunk
        if isinstance(obj, DocumentChunk):
            added_chunks.append(obj.content)

    assert len(added_chunks) > 0
    
    # The chunk text must contain BOTH the original slide text and the LLM image description
    combined_chunk_text = " ".join(added_chunks)
    assert "Slide text description." in combined_chunk_text
    assert "[Slide Image 1 Description]" in combined_chunk_text
    assert "A red circle chart." in combined_chunk_text

